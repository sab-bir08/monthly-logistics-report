"""
Monthly Logistics Performance Report  (synthetic-data demo)
===========================================================

Turns raw operational logistics data into an executive-ready deck - the same
workflow used for a real (confidential) monthly report, here on SYNTHETIC data.

It does two things:
  1. Renders three clean SVG charts (used in the README).
  2. Builds a real, editable PowerPoint deck with NATIVE charts + a summary
     table using python-pptx (logistics_report.pptx).

Run:
    pip install -r requirements.txt
    python build_report.py

Author: Sabbir | github.com/sab-bir08
"""
from __future__ import annotations
import calendar
import datetime as dt
import random
from pathlib import Path
from xml.sax.saxutils import escape as esc

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
SEED = 8
INK, GRID, MUTED = "#1f2937", "#e5e7eb", "#6b7280"
PALETTE = ["#2563eb", "#0891b2", "#16a34a", "#f59e0b", "#dc2626"]
CATS = ["Standard", "Express", "International", "Returns"]


# ---------------------------------------------------------------------------
# Synthetic operational data
# ---------------------------------------------------------------------------
def make_data():
    rng = random.Random(SEED)
    today = dt.date.today().replace(day=1)
    month = today.replace(day=1) - dt.timedelta(days=1)   # previous full month
    ndays = calendar.monthrange(month.year, month.month)[1]
    days, dispatches = [], []
    for d in range(1, ndays + 1):
        date = dt.date(month.year, month.month, d)
        weekday = date.weekday()
        base = 820 if weekday < 5 else 280            # weekend dip
        val = int(base * (1 + rng.gauss(0, 0.10)) + d * 4)  # slight upward drift
        days.append(d)
        dispatches.append(max(0, val))
    total = sum(dispatches)
    cat_share = [0.52, 0.24, 0.14, 0.10]
    cats = {c: int(total * s) for c, s in zip(CATS, cat_share)}
    # back-office: processed vs backlog per week
    weeks = ["W1", "W2", "W3", "W4", "W5"][:((ndays + 6) // 7)]
    processed, backlog = [], []
    for _ in weeks:
        processed.append(int(rng.uniform(3600, 5200)))
        backlog.append(int(rng.uniform(180, 760)))
    return {"month": month.strftime("%B %Y"), "days": days, "dispatches": dispatches,
            "cats": cats, "weeks": weeks, "processed": processed, "backlog": backlog,
            "total": total}


# ---------------------------------------------------------------------------
# SVG helpers
# ---------------------------------------------------------------------------
def _hdr(w, h):
    return (f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {w} {h}" '
            f'font-family="Segoe UI, Helvetica, Arial, sans-serif">'
            f'<rect width="{w}" height="{h}" fill="#ffffff"/>')


def line_svg(days, vals, title):
    w, h, ml, mr, mt, mb = 720, 300, 48, 20, 48, 40
    pw, ph = w - ml - mr, h - mt - mb
    vmax = max(vals) * 1.12
    X = lambda i: ml + pw * i / (len(vals) - 1)
    Y = lambda v: mt + ph * (1 - v / vmax)
    s = [_hdr(w, h)]
    s.append(f'<text x="{ml}" y="28" font-size="16" font-weight="700" fill="{INK}">{esc(title)}</text>')
    for g in range(5):
        gy = mt + ph * g / 4
        s.append(f'<line x1="{ml}" y1="{gy:.0f}" x2="{w-mr}" y2="{gy:.0f}" stroke="{GRID}"/>')
    pts = " ".join(f"{X(i):.1f},{Y(v):.1f}" for i, v in enumerate(vals))
    s.append(f'<polygon points="{ml},{mt+ph} {pts} {w-mr},{mt+ph}" fill="{PALETTE[0]}" opacity="0.08"/>')
    s.append(f'<polyline points="{pts}" fill="none" stroke="{PALETTE[0]}" stroke-width="2.4"/>')
    for i, d in enumerate(days):
        if d % 5 == 0 or d == 1:
            s.append(f'<text x="{X(i):.1f}" y="{h-14}" font-size="10" text-anchor="middle" fill="{MUTED}">{d}</text>')
    s.append(f'<text x="{ml}" y="{h-14}" font-size="10" fill="{MUTED}">day of month</text>')
    s.append('</svg>')
    return "\n".join(s)


def donut_svg(cats, title):
    import math
    w, h = 720, 300
    cx, cy, r, rin = 170, 165, 95, 55
    total = sum(cats.values())
    s = [_hdr(w, h)]
    s.append(f'<text x="24" y="28" font-size="16" font-weight="700" fill="{INK}">{esc(title)}</text>')
    a0 = -math.pi / 2
    for i, (c, v) in enumerate(cats.items()):
        frac = v / total
        a1 = a0 + frac * 2 * math.pi
        large = 1 if frac > 0.5 else 0
        x0, y0 = cx + r * math.cos(a0), cy + r * math.sin(a0)
        x1, y1 = cx + r * math.cos(a1), cy + r * math.sin(a1)
        xi0, yi0 = cx + rin * math.cos(a1), cy + rin * math.sin(a1)
        xi1, yi1 = cx + rin * math.cos(a0), cy + rin * math.sin(a0)
        col = PALETTE[i % len(PALETTE)]
        s.append(f'<path d="M{x0:.1f},{y0:.1f} A{r},{r} 0 {large} 1 {x1:.1f},{y1:.1f} '
                 f'L{xi0:.1f},{yi0:.1f} A{rin},{rin} 0 {large} 0 {xi1:.1f},{yi1:.1f} Z" fill="{col}"/>')
        a0 = a1
    s.append(f'<text x="{cx}" y="{cy-2}" font-size="20" font-weight="700" text-anchor="middle" fill="{INK}">{total/1000:.1f}k</text>')
    s.append(f'<text x="{cx}" y="{cy+16}" font-size="10" text-anchor="middle" fill="{MUTED}">dispatches</text>')
    ly = 70
    for i, (c, v) in enumerate(cats.items()):
        col = PALETTE[i % len(PALETTE)]
        s.append(f'<rect x="320" y="{ly}" width="14" height="14" rx="3" fill="{col}"/>')
        s.append(f'<text x="342" y="{ly+12}" font-size="12" fill="{INK}">{esc(c)}</text>')
        s.append(f'<text x="690" y="{ly+12}" font-size="12" text-anchor="end" fill="{MUTED}">{v:,} ({100*v/total:.0f}%)</text>')
        ly += 34
    s.append('</svg>')
    return "\n".join(s)


def grouped_bar_svg(weeks, processed, backlog, title):
    w, h, ml, mr, mt, mb = 720, 300, 52, 20, 56, 40
    pw, ph = w - ml - mr, h - mt - mb
    vmax = max(processed) * 1.15
    n = len(weeks)
    gw = pw / n
    s = [_hdr(w, h)]
    s.append(f'<text x="{ml}" y="28" font-size="16" font-weight="700" fill="{INK}">{esc(title)}</text>')
    for g in range(5):
        gy = mt + ph * g / 4
        s.append(f'<line x1="{ml}" y1="{gy:.0f}" x2="{w-mr}" y2="{gy:.0f}" stroke="{GRID}"/>')
    bw = gw * 0.32
    for i, wk in enumerate(weeks):
        cxv = ml + gw * i + gw / 2
        hp = ph * processed[i] / vmax
        hb = ph * backlog[i] / vmax
        s.append(f'<rect x="{cxv-bw-2:.1f}" y="{mt+ph-hp:.1f}" width="{bw:.1f}" height="{hp:.1f}" rx="3" fill="{PALETTE[2]}"/>')
        s.append(f'<rect x="{cxv+2:.1f}" y="{mt+ph-hb:.1f}" width="{bw:.1f}" height="{hb:.1f}" rx="3" fill="{PALETTE[3]}"/>')
        s.append(f'<text x="{cxv:.1f}" y="{h-14}" font-size="11" text-anchor="middle" fill="{MUTED}">{esc(wk)}</text>')
    s.append(f'<rect x="{w-220}" y="40" width="12" height="12" rx="2" fill="{PALETTE[2]}"/>')
    s.append(f'<text x="{w-202}" y="50" font-size="11" fill="{INK}">Processed</text>')
    s.append(f'<rect x="{w-120}" y="40" width="12" height="12" rx="2" fill="{PALETTE[3]}"/>')
    s.append(f'<text x="{w-102}" y="50" font-size="11" fill="{INK}">Backlog</text>')
    s.append('</svg>')
    return "\n".join(s)


# ---------------------------------------------------------------------------
# PowerPoint deck (native, editable charts)
# ---------------------------------------------------------------------------
def build_pptx(data):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank, titleonly = prs.slide_layouts[6], prs.slide_layouts[5]

    # title slide
    s = prs.slides.add_slide(blank)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2))
    tf = box.text_frame
    tf.text = "Monthly Logistics Performance Report"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = f"{data['month']}  -  Operations review"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    def chart_slide(title, ctype, cats, series, legend=False):
        sl = prs.slides.add_slide(titleonly)
        sl.shapes.title.text = title
        cd = CategoryChartData()
        cd.categories = cats
        for name, vals in series:
            cd.add_series(name, vals)
        gf = sl.shapes.add_chart(ctype, Inches(0.8), Inches(1.6),
                                 Inches(11.7), Inches(5.3), cd)
        ch = gf.chart
        ch.has_legend = legend
        if legend:
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
        return ch

    chart_slide("Daily dispatches", XL_CHART_TYPE.LINE,
                [str(d) for d in data["days"]],
                [("Dispatches", data["dispatches"])])
    chart_slide("Dispatch volume by category", XL_CHART_TYPE.DOUGHNUT,
                list(data["cats"].keys()),
                [("Volume", list(data["cats"].values()))], legend=True)
    chart_slide("Back-office throughput", XL_CHART_TYPE.COLUMN_CLUSTERED,
                data["weeks"],
                [("Processed", data["processed"]), ("Backlog", data["backlog"])],
                legend=True)

    # summary table slide
    sl = prs.slides.add_slide(titleonly)
    sl.shapes.title.text = "Summary"
    rows = [("Total dispatches", f"{data['total']:,}"),
            ("Busiest category", f"{max(data['cats'], key=data['cats'].get)}"),
            ("Avg daily (weekday)", f"{int(sum(data['dispatches'])/len(data['days'])):,}"),
            ("Total processed (back office)", f"{sum(data['processed']):,}"),
            ("Closing backlog", f"{data['backlog'][-1]:,}")]
    tbl = sl.shapes.add_table(len(rows) + 1, 2, Inches(1.5), Inches(1.8),
                              Inches(10), Inches(4)).table
    tbl.cell(0, 0).text, tbl.cell(0, 1).text = "Metric", "Value"
    for i, (k, v) in enumerate(rows, 1):
        tbl.cell(i, 0).text, tbl.cell(i, 1).text = k, v
    out = ROOT / "logistics_report.pptx"
    prs.save(out)
    return out


def main():
    data = make_data()
    ASSETS.mkdir(exist_ok=True)
    (ASSETS / "daily_dispatches.svg").write_text(
        line_svg(data["days"], data["dispatches"], f"Daily dispatches - {data['month']}"), encoding="utf-8")
    (ASSETS / "category_breakdown.svg").write_text(
        donut_svg(data["cats"], "Dispatch volume by category"), encoding="utf-8")
    (ASSETS / "throughput.svg").write_text(
        grouped_bar_svg(data["weeks"], data["processed"], data["backlog"],
                        "Back-office throughput by week"), encoding="utf-8")
    print(f"Logistics report - {data['month']}")
    print(f"  total dispatches: {data['total']:,}")
    print(f"  busiest category: {max(data['cats'], key=data['cats'].get)}")
    print("  wrote 3 SVG charts to assets/")
    try:
        out = build_pptx(data)
        print(f"  built editable deck: {out.name}")
    except Exception as e:                       # python-pptx optional
        print(f"  (skipped .pptx: {e})")


if __name__ == "__main__":
    main()
